from __future__ import annotations
import tempfile
from pathlib import Path
from typing import Optional
import fitz
from docx import Document
from openpyxl import Workbook, load_workbook
from openpyxl.chart import BarChart, Reference
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches
from app.conversion_engine import SPECS, validate_output_file
from app.processing_quality import run_conversion

def _run(tool_id: str, source: Path, output: Path, root: Path, options: Optional[dict]=None) -> Path:
    spec = SPECS[tool_id]
    workdir = root / f'work-{tool_id}-{output.stem}'
    workdir.mkdir(parents=True, exist_ok=True)
    run_conversion(spec, [source], output, options or {}, workdir, None)
    validate_output_file(output, spec.output_extension)
    return output

def _docx_fixture(path: Path, image_path: Path) -> None:
    doc = Document()
    section = doc.sections[0]
    section.header.paragraphs[0].text = 'AJN R19 HEADER'
    section.footer.paragraphs[0].text = 'AJN R19 FOOTER'
    doc.add_heading('AJN PDF R19 Word Fidelity', 0)
    doc.add_paragraph('Word conversion must preserve readable text, tables, images and pagination intent.')
    table = doc.add_table(rows=3, cols=3)
    for r, values in enumerate([['Product', 'Count', 'Amount'], ['AJN PDF', '12', '480']]):
        for c, value in enumerate(values):
            table.cell(r, c).text = value
    doc.add_picture(str(image_path), width=Inches(2.0))
    doc.add_paragraph('END WORD MARKER R19')
    doc.save(path)

def _xlsx_fixture(path: Path) -> None:
    wb = Workbook()
    ws = wb.active
    ws.title = 'Summary'
    ws.append(['Product', 'Count', 'Amount'])
    ws.append(['AJN PDF', 12, 480])
    ws.append([6, 240])
    ws['E1'] = 'AJN R19 EXCEL MARKER'
    ws.merge_cells('E1:G1')
    ws['D2'] = '=B2*C2'
    chart = BarChart()
    data = Reference(ws, min_col=3, min_row=1, max_row=3)
    cats = Reference(ws, min_col=1, min_row=2, max_row=3)
    chart.add_data(data, titles_from_data=True)
    chart.set_categories(cats)
    ws.add_chart(chart, 'E4')
    ws.page_setup.orientation = 'landscape'
    details = wb.create_sheet('Details')
    details.append(['ID', 'Name'])
    details.append([1, 'Second Sheet R19'])
    wb.save(path)

def _pptx_fixture(path: Path, image_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for title, body in [('AJN R19 SLIDE ONE', 'Presentation fidelity marker one'), ('AJN R19 SLIDE TWO', 'Presentation fidelity marker two')]:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = body
        slide.shapes.add_picture(str(image_path), Inches(8.5), Inches(4.2), width=Inches(2))
    prs.save(path)

def _pdf_fixture(path: Path) -> None:
    doc = fitz.open()
    page = doc.new_page(width=595, height=842)
    page.insert_text((60, 70), 'AJN PDF R19 PORTRAIT DOCUMENT', fontsize=20)
    page.insert_text((60, 105), 'Editable PowerPoint marker and Word layout marker.', fontsize=11)
    x0, y0 = (60, 180)
    widths = [180, 120, 120]
    row_h = 34
    rows = [['Product', 'Count', 'Amount'], ['AJN PDF', '12', '480']]
    xs = [x0]
    for w in widths:
        xs.append(xs[-1] + w)
    for ri in range(len(rows) + 1):
        page.draw_line((x0, y0 + ri * row_h), (xs[-1], y0 + ri * row_h))
    for x in xs:
        page.draw_line((x, y0), (x, y0 + len(rows) * row_h))
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            page.insert_text((xs[ci] + 7, y0 + ri * row_h + 22), val, fontsize=10)
    page2 = doc.new_page(width=842, height=595)
    page2.insert_text((70, 80), 'AJN R19 LANDSCAPE SECOND PAGE', fontsize=18)
    doc.save(path)
    doc.close()

def main() -> None:
    with tempfile.TemporaryDirectory(prefix='ajn-r19-fidelity-') as temporary:
        root = Path(temporary)
        image_path = root / 'logo.png'
        image = Image.new('RGB', (500, 220), 'white')
        draw = ImageDraw.Draw(image)
        draw.rectangle((20, 20, 480, 200), outline='black', width=4)
        draw.text((70, 90), 'AJN PDF R19', fill='black')
        image.save(image_path)
        image.close()
        docx = root / 'complex.docx'
        xlsx = root / 'complex.xlsx'
        pptx = root / 'complex.pptx'
        pdf = root / 'source.pdf'
        _docx_fixture(docx, image_path)
        _xlsx_fixture(xlsx)
        _pptx_fixture(pptx, image_path)
        _pdf_fixture(pdf)
        word_pdf = _run('word-to-pdf', docx, root / 'word.pdf', root, {'office_fit_mode': 'preserve'})
        with fitz.open(word_pdf) as out:
            text = ''.join((p.get_text() for p in out))
            assert 'AJN PDF R19 Word Fidelity' in text
            assert 'END WORD MARKER R19' in text
        excel_pdf = _run('excel-to-pdf', xlsx, root / 'excel.pdf', root, {'office_fit_mode': 'preserve'})
        with fitz.open(excel_pdf) as out:
            text = ''.join((p.get_text() for p in out))
            assert 'AJN' in text and ('480' in text or '240' in text)
        ppt_pdf = _run('powerpoint-to-pdf', pptx, root / 'ppt.pdf', root, {'export_notes': False})
        with fitz.open(ppt_pdf) as out:
            assert out.page_count == 2
            text = ''.join((p.get_text() for p in out))
            assert 'AJN R19 SLIDE ONE' in text
            assert 'AJN R19 SLIDE TWO' in text
        preserve = _run('pdf-to-pptx', pdf, root / 'preserve.pptx', root, {'ppt_mode': 'preserve', 'dpi': 150})
        prs = Presentation(preserve)
        assert len(prs.slides) == 2
        source_ratio = 595 / 842
        slide_ratio = prs.slide_width / prs.slide_height
        assert abs(source_ratio - slide_ratio) < 0.03, (source_ratio, slide_ratio)
        assert any((shape.shape_type == 13 for shape in prs.slides[0].shapes)), 'Preserve mode must contain a page image'
        editable = _run('pdf-to-pptx', pdf, root / 'editable.pptx', root, {'ppt_mode': 'editable'})
        eprs = Presentation(editable)
        editable_text = ''.join((shape.text for slide in eprs.slides for shape in slide.shapes if hasattr(shape, 'text')))
        assert 'AJN PDF R19 PORTRAIT DOCUMENT' in editable_text
        word = _run('pdf-to-word', pdf, root / 'layout.docx', root, {'include_images': True, 'language': 'eng', 'dpi': 240})
        word_doc = Document(word)
        word_text = ''.join((paragraph.text for paragraph in word_doc.paragraphs))
        assert 'AJN PDF R19 PORTRAIT DOCUMENT' in word_text
        table_text = '|'.join((cell.text for table in word_doc.tables for row in table.rows for cell in row.cells))
        assert 'Product' in table_text and '480' in table_text
        excel = _run('pdf-to-excel', pdf, root / 'table.xlsx', root, {})
        wb = load_workbook(excel, data_only=False)
        try:
            joined = '|'.join((str(cell.value or '') for ws in wb.worksheets for row in ws.iter_rows() for cell in row))
            assert 'Product' in joined and 'AJN PDF' in joined and ('480' in joined)
        finally:
            wb.close()
        csv_out = _run('pdf-to-csv', pdf, root / 'table.csv', root, {})
        csv_text = csv_out.read_text(encoding='utf-8-sig')
        assert 'Product' in csv_text and '480' in csv_text
        print('PASS: AJN PDF R19 conversion fidelity — Word/Excel/PPT to PDF, PDF to Word layout/table, PDF to PPT preserve/editable, PDF table to XLSX/CSV')
if __name__ == '__main__':
    main()
