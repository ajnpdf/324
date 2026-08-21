from __future__ import annotations
import csv
import json
import math
import re
import shutil
import statistics
from pathlib import Path
from typing import Any, Optional
import fitz
from openpyxl import Workbook, load_workbook
from openpyxl.styles import Alignment, Font, PatternFill
from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt
from . import conversion_engine as legacy
WRITER_EXTENSIONS = {'.doc', '.docx', '.odt', '.rtf'}
CALC_EXTENSIONS = {'.xls', '.xlsx', '.ods'}
IMPRESS_EXTENSIONS = {'.ppt', '.pptx', '.odp'}

def _lo_filter(source: Path, options: dict[str, Any]) -> str:
    ext = source.suffix.lower()
    if ext in WRITER_EXTENSIONS:
        return 'pdf:writer_pdf_Export'
    if ext in CALC_EXTENSIONS:
        if str(options.get('office_fit_mode', 'preserve')).lower() == 'single-page-sheet':
            data = json.dumps({'SinglePageSheets': {'type': 'boolean', 'value': 'true'}}, separators=(',', ':'))
            return f'pdf:calc_pdf_Export:{data}'
        return 'pdf:calc_pdf_Export'
    if ext in IMPRESS_EXTENSIONS:
        if bool(options.get('export_notes', False)):
            data = json.dumps({'ExportNotesPages': {'type': 'boolean', 'value': 'true'}}, separators=(',', ':'))
            return f'pdf:impress_pdf_Export:{data}'
        return 'pdf:impress_pdf_Export'
    return 'pdf'

def _source_markers(source: Path) -> tuple[list[str], Optional[int]]:
    markers: list[str] = []
    expected_pages: Optional[int] = None
    ext = source.suffix.lower()
    try:
        if ext == '.docx':
            from docx import Document
            doc = Document(source)
            blocks = [p.text for p in doc.paragraphs]
            blocks += [cell.text for table in doc.tables for row in table.rows for cell in row.cells]
            for text in blocks:
                markers.extend(re.findall('[A-Za-z0-9][A-Za-z0-9._-]{2,}', text or '')[:4])
        elif ext == '.xlsx':
            wb = load_workbook(source, read_only=True, data_only=False)
            try:
                visible = [ws for ws in wb.worksheets if ws.sheet_state == 'visible']
                expected_pages = max(1, len(visible))
                for ws in visible[:5]:
                    markers.append(ws.title)
                    used = 0
                    for row in ws.iter_rows(values_only=True):
                        for value in row:
                            if value not in (None, ''):
                                markers.extend(re.findall('[A-Za-z0-9][A-Za-z0-9._-]{2,}', str(value))[:2])
                                used += 1
                                if used >= 20:
                                    break
                        if used >= 20:
                            break
            finally:
                wb.close()
        elif ext == '.pptx':
            prs = Presentation(source)
            expected_pages = len(prs.slides)
            for slide in prs.slides[:8]:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        markers.extend(re.findall('[A-Za-z0-9][A-Za-z0-9._-]{2,}', str(shape.text or ''))[:4])
    except Exception:
        return ([], expected_pages)
    unique: list[str] = []
    seen: set[str] = set()
    for marker in markers:
        key = marker.casefold()
        if len(marker) >= 3 and key not in seen:
            seen.add(key)
            unique.append(marker)
        if len(unique) >= 20:
            break
    return (unique, expected_pages)

def _verify_office_pdf(source: Path, output: Path) -> None:
    markers, expected_pages = _source_markers(source)
    with fitz.open(output) as pdf:
        if pdf.page_count < 1:
            raise RuntimeError('Office converter created a PDF with no pages.')
        if source.suffix.lower() == '.pptx' and expected_pages is not None and (pdf.page_count != expected_pages):
            raise RuntimeError(f'PowerPoint fidelity check failed: expected{expected_pages}pages, got{pdf.page_count}.')
        searchable = ''.join((page.get_text('text') for page in pdf))
    if markers and searchable.strip():
        lower = searchable.casefold()
        hits = sum((1 for marker in markers if marker.casefold() in lower))
        required = max(1, min(3, math.ceil(len(markers) * 0.2)))
        if hits < required:
            raise RuntimeError(f'Office fidelity check failed: only{hits}/{len(markers)}sampled markers were preserved.')

def office_to_pdf_fidelity(source: Path, output: Path, workdir: Path, options: dict[str, Any]) -> None:
    executable = legacy.command_path('libreoffice')
    if not executable:
        raise RuntimeError('LibreOffice is not installed on the processing server.')
    last_error: Optional[Exception] = None
    for attempt in range(1, 3):
        profile = workdir / f'r19-lo-profile-{attempt}'
        shutil.rmtree(profile, ignore_errors=True)
        profile.mkdir(parents=True, exist_ok=True)
        generated = workdir / f'{source.stem}.pdf'
        generated.unlink(missing_ok=True)
        command = [executable, '--headless', '--nologo', '--nodefault', '--nolockcheck', '--norestore', f'-env:UserInstallation={profile.resolve().as_uri()}', '--convert-to', _lo_filter(source, options), '--outdir', str(workdir), str(source)]
        try:
            legacy._run(command, timeout=180)
            if not generated.exists() or generated.stat().st_size <= 32:
                candidates = [p for p in workdir.glob('*.pdf') if p != output and p.stat().st_size > 32]
                if not candidates:
                    raise RuntimeError('LibreOffice did not create a PDF output.')
                generated = max(candidates, key=lambda p: p.stat().st_mtime)
            shutil.move(str(generated), str(output))
            legacy.validate_output_file(output, '.pdf')
            _verify_office_pdf(source, output)
            return
        except Exception as exc:
            last_error = exc
            output.unlink(missing_ok=True)
        finally:
            shutil.rmtree(profile, ignore_errors=True)
    raise RuntimeError(f'LibreOffice fidelity conversion failed after retry:{last_error}')

def _fit_rect(container_w: int, container_h: int, content_w: float, content_h: float) -> tuple[int, int, int, int]:
    scale = min(container_w / max(1.0, content_w), container_h / max(1.0, content_h))
    width, height = (max(1, int(content_w * scale)), max(1, int(content_h * scale)))
    return ((container_w - width) // 2, (container_h - height) // 2, width, height)

def _span_rgb(value: int) -> RGBColor:
    return RGBColor(value >> 16 & 255, value >> 8 & 255, value & 255)

def pdf_to_pptx_fidelity(source: Path, output: Path, options: dict[str, Any], workdir: Path) -> None:
    mode = str(options.get('ppt_mode', 'preserve')).lower()
    if mode not in {'preserve', 'editable'}:
        raise ValueError('PowerPoint mode must be preserve or editable.')
    dpi = max(96, min(300, int(options.get('dpi', 160))))
    prs = Presentation()
    blank = prs.slide_layouts[6]
    with fitz.open(source) as pdf:
        if pdf.page_count < 1:
            raise ValueError('The PDF has no pages.')
        first = pdf[0].rect
        prs.slide_width = int(PptInches(first.width / 72.0))
        prs.slide_height = int(PptInches(first.height / 72.0))
        for page_number, page in enumerate(pdf, start=1):
            slide = prs.slides.add_slide(blank)
            if mode == 'preserve':
                pix = page.get_pixmap(matrix=fitz.Matrix(dpi / 72.0, dpi / 72.0), alpha=False)
                image_path = workdir / f'r19-slide-{page_number:03d}.png'
                pix.save(str(image_path))
                left, top, width, height = _fit_rect(prs.slide_width, prs.slide_height, pix.width, pix.height)
                slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
                image_path.unlink(missing_ok=True)
                continue
            sx = prs.slide_width / max(1.0, page.rect.width)
            sy = prs.slide_height / max(1.0, page.rect.height)
            for block in page.get_text('dict', sort=True).get('blocks', []):
                if int(block.get('type', 0)) != 0:
                    continue
                x0, y0, x1, y1 = [float(v) for v in block.get('bbox', (0, 0, 0, 0))]
                if x1 <= x0 or y1 <= y0:
                    continue
                box = slide.shapes.add_textbox(int(x0 * sx), int(y0 * sy), max(1, int((x1 - x0) * sx)), max(1, int((y1 - y0) * sy)))
                frame = box.text_frame
                frame.clear()
                frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0
                first_paragraph = True
                for line in block.get('lines', []):
                    paragraph = frame.paragraphs[0] if first_paragraph else frame.add_paragraph()
                    first_paragraph = False
                    for span in line.get('spans', []):
                        text = str(span.get('text') or '')
                        if not text:
                            continue
                        run = paragraph.add_run()
                        run.text = text
                        run.font.size = PptPt(max(6.0, min(60.0, float(span.get('size') or 11))))
                        font = str(span.get('font') or '').lower()
                        flags = int(span.get('flags') or 0)
                        run.font.bold = bool(flags & 16) or 'bold' in font
                        run.font.italic = bool(flags & 2) or 'italic' in font or 'oblique' in font
                        run.font.color.rgb = _span_rgb(int(span.get('color') or 0))
            seen: set[int] = set()
            for info in page.get_images(full=True):
                xref = int(info[0])
                if xref in seen:
                    continue
                seen.add(xref)
                try:
                    extracted = pdf.extract_image(xref)
                    data = extracted.get('image')
                    if not data:
                        continue
                    temp = workdir / f"r19-image-{page_number}-{xref}.{extracted.get('ext') or 'png'}"
                    temp.write_bytes(data)
                    for rect in page.get_image_rects(xref)[:4]:
                        slide.shapes.add_picture(str(temp), int(rect.x0 * sx), int(rect.y0 * sy), width=max(1, int(rect.width * sx)), height=max(1, int(rect.height * sy)))
                    temp.unlink(missing_ok=True)
                except Exception:
                    continue
    if len(prs.slides) < 1:
        raise RuntimeError('PDF to PowerPoint created no slides.')
    prs.save(output)

def _bbox_inside(inner: tuple[float, float, float, float], outer: tuple[float, float, float, float]) -> bool:
    ix0, iy0, ix1, iy1 = inner
    ox0, oy0, ox1, oy1 = outer
    cx, cy = ((ix0 + ix1) / 2.0, (iy0 + iy1) / 2.0)
    return ox0 <= cx <= ox1 and oy0 <= cy <= oy1

def pdf_to_docx_fidelity(source: Path, output: Path, options: dict[str, Any], workdir: Path) -> None:
    include_images = bool(options.get('include_images', True))
    document = Document()
    document.core_properties.title = source.stem
    with fitz.open(source) as pdf:
        if pdf.page_count < 1:
            raise ValueError('The PDF has no pages.')
        for page_index, page in enumerate(pdf):
            if page_index == 0:
                section = document.sections[0]
            else:
                section = document.add_section(WD_SECTION.NEW_PAGE)
            section.page_width = Pt(float(page.rect.width))
            section.page_height = Pt(float(page.rect.height))
            section.top_margin = section.bottom_margin = Pt(28)
            section.left_margin = section.right_margin = Pt(28)
            raw = page.get_text('dict', sort=True)
            native_text = page.get_text('text', sort=True).strip()
            if len(re.sub('\\s+', '', native_text)) < 24:
                recognized_text = ''
                if not recognized_text:
                    document.add_paragraph('[No readable text detected on this page]')
                else:
                    for paragraph_text in re.split('\\n\\s*\\n', recognized_text):
                        if paragraph_text.strip():
                            document.add_paragraph(paragraph_text.strip())
                continue
            tables: list[dict[str, Any]] = []
            try:
                for table in list(getattr(page.find_tables(), 'tables', []) or []):
                    bbox = tuple((float(v) for v in table.bbox))
                    rows = [[str(v or '').strip() for v in row] for row in table.extract() or []]
                    rows = [row for row in rows if any(row)]
                    if rows:
                        tables.append({'kind': 'table', 'bbox': bbox, 'rows': rows})
            except Exception:
                tables = []
            items: list[dict[str, Any]] = list(tables)
            for block in raw.get('blocks', []):
                if int(block.get('type', 0)) != 0:
                    continue
                bbox = tuple((float(v) for v in block.get('bbox', (0, 0, 0, 0))))
                if any((_bbox_inside(bbox, table['bbox']) for table in tables)):
                    continue
                items.append({'kind': 'text', 'bbox': bbox, 'block': block})
            if include_images:
                try:
                    for info in page.get_image_info(xrefs=True):
                        xref = int(info.get('xref') or 0)
                        bbox = tuple((float(v) for v in info.get('bbox', (0, 0, 0, 0))))
                        if xref > 0 and bbox[2] > bbox[0] and (bbox[3] > bbox[1]):
                            items.append({'kind': 'image', 'bbox': bbox, 'xref': xref})
                except Exception:
                    pass
            items.sort(key=lambda item: (item['bbox'][1], item['bbox'][0]))
            font_sizes = [float(span.get('size') or 11) for block in raw.get('blocks', []) if int(block.get('type', 0)) == 0 for line in block.get('lines', []) for span in line.get('spans', []) if 5 <= float(span.get('size') or 11) <= 48]
            body_size = statistics.median(font_sizes) if font_sizes else 11.0
            for item_index, item in enumerate(items):
                if item['kind'] == 'table':
                    rows = item['rows']
                    cols = max((len(row) for row in rows))
                    table = document.add_table(rows=len(rows), cols=cols)
                    table.style = 'Table Grid'
                    for ri, row in enumerate(rows):
                        for ci, value in enumerate(row):
                            cell = table.cell(ri, ci)
                            cell.text = value
                            if ri == 0:
                                for run in cell.paragraphs[0].runs:
                                    run.bold = True
                    continue
                if item['kind'] == 'image':
                    try:
                        extracted = pdf.extract_image(item['xref'])
                        data = extracted.get('image')
                        if not data or len(data) < 512:
                            continue
                        suffix = extracted.get('ext') or 'png'
                        image_path = workdir / f'r19-docx-image-{page_index + 1}-{item_index}.{suffix}'
                        image_path.write_bytes(data)
                        width_inches = max(0.5, min(6.5, (item['bbox'][2] - item['bbox'][0]) / 72.0))
                        document.add_picture(str(image_path), width=Inches(width_inches))
                        image_path.unlink(missing_ok=True)
                    except Exception:
                        continue
                    continue
                block = item['block']
                spans = [span for line in block.get('lines', []) for span in line.get('spans', []) if str(span.get('text') or '')]
                if not spans:
                    continue
                average = statistics.mean((float(span.get('size') or body_size) for span in spans))
                style = 'Heading 1' if average >= body_size * 1.65 else 'Heading 2' if average >= body_size * 1.3 else None
                paragraph = document.add_paragraph(style=style)
                x0, _, x1, _ = item['bbox']
                paragraph.paragraph_format.left_indent = Pt(max(0, x0 - 28))
                paragraph.paragraph_format.right_indent = Pt(max(0, float(page.rect.width) - x1 - 28))
                center = (x0 + x1) / 2.0
                if abs(center - float(page.rect.width) / 2.0) < float(page.rect.width) * 0.08 and x1 - x0 < float(page.rect.width) * 0.8:
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                first_line = True
                for line in block.get('lines', []):
                    if not first_line:
                        paragraph.add_run().add_break()
                    first_line = False
                    for span in line.get('spans', []):
                        text = str(span.get('text') or '')
                        if not text:
                            continue
                        run = paragraph.add_run(text)
                        font_name = str(span.get('font') or '').lower()
                        flags = int(span.get('flags') or 0)
                        run.bold = bool(flags & 16) or 'bold' in font_name
                        run.italic = bool(flags & 2) or 'italic' in font_name or 'oblique' in font_name
                        run.font.size = Pt(max(6.0, min(48.0, float(span.get('size') or body_size))))
    if document.paragraphs and (not document.paragraphs[0].text) and (len(document.paragraphs) > 1):
        element = document.paragraphs[0]._element
        element.getparent().remove(element)
    document.save(output)

def _coerce(value: str) -> Any:
    text = str(value or '').strip()
    normalized = text.replace(',', '')
    if re.fullmatch('[-+]?\\d+', normalized):
        try:
            return int(normalized)
        except Exception:
            pass
    if re.fullmatch('[-+]?\\d*\\.\\d+', normalized):
        try:
            return float(normalized)
        except Exception:
            pass
    currency = re.fullmatch('[₹$€£]\\s*([-+]?\\d+(?:,\\d{3})*(?:\\.\\d+)?)', text)
    if currency:
        try:
            return float(currency.group(1).replace(',', ''))
        except Exception:
            pass
    return text

def _detected_tables(source: Path) -> list[tuple[int, int, list[list[str]]]]:
    found: list[tuple[int, int, list[list[str]]]] = []
    with fitz.open(source) as pdf:
        for page_number, page in enumerate(pdf, start=1):
            try:
                tables = list(getattr(page.find_tables(), 'tables', []) or [])
            except Exception:
                tables = []
            for table_number, table in enumerate(tables, start=1):
                rows = [[str(v or '').strip() for v in row] for row in table.extract() or []]
                rows = [row for row in rows if any(row)]
                if rows and max((len(r) for r in rows), default=0) >= 2:
                    found.append((page_number, table_number, rows))
    return found

def _write_xlsx(tables: list[tuple[int, int, list[list[str]]]], output: Path) -> None:
    wb = Workbook()
    wb.remove(wb.active)
    for page_number, table_number, rows in tables:
        ws = wb.create_sheet(f'P{page_number}-Table{table_number}'[:31])
        for ri, row in enumerate(rows, start=1):
            for ci, value in enumerate(row, start=1):
                cell = ws.cell(ri, ci, _coerce(value))
                cell.alignment = Alignment(vertical='top', wrap_text=True)
                if ri == 1:
                    cell.font = Font(bold=True)
                    cell.fill = PatternFill('solid', fgColor='E8EEF8')
        for column in ws.columns:
            ws.column_dimensions[column[0].column_letter].width = min(48, max(10, max((len(str(c.value or '')) for c in column)) + 2))
        ws.freeze_panes = 'A2'
    if not wb.sheetnames:
        raise ValueError('No structured table was detected in this PDF.')
    wb.save(output)

def pdf_to_xlsx_fidelity(source: Path, output: Path, options: dict[str, Any]) -> None:
    tables = _detected_tables(source)
    if not tables and bool(options.get('scanned_table_fallback', False)):
        tables = []
    if not tables:
        if bool(options.get('allow_unstructured', False)):
            wb = Workbook()
            ws = wb.active
            ws.title = 'PDF text'
            ws.append(['Page', 'Line', 'Text'])
            with fitz.open(source) as pdf:
                for page_number, page in enumerate(pdf, start=1):
                    for line_number, line in enumerate(page.get_text('text', sort=True).splitlines(), start=1):
                        if line.strip():
                            ws.append([page_number, line_number, line.strip()])
            wb.save(output)
            return
        raise ValueError('No structured table was detected in this PDF. Enable scanned-table fallback for scanned tables or use PDF to Word/Text.')
    _write_xlsx(tables, output)

def pdf_to_csv_fidelity(source: Path, output: Path, options: dict[str, Any]) -> None:
    tables = _detected_tables(source)
    if not tables and bool(options.get('scanned_table_fallback', False)):
        tables = []
    if not tables:
        raise ValueError('No structured table was detected in this PDF.')
    with output.open('w', encoding='utf-8-sig', newline='') as handle:
        writer = csv.writer(handle)
        for index, (page_number, table_number, rows) in enumerate(tables):
            if index:
                writer.writerow([])
            writer.writerow([f'Page{page_number}- Table{table_number}'])
            writer.writerows(rows)

def run_r19_conversion(spec: legacy.ConversionSpec, files: list[Path], output: Path, options: dict[str, Any], workdir: Path, source_url: Optional[str]=None) -> bool:
    del source_url
    source = files[0] if files else None
    if source is None:
        return False
    if spec.processor == 'office_to_pdf':
        office_to_pdf_fidelity(source, output, workdir, options)
        return True
    if spec.processor == 'pdf_docx':
        pdf_to_docx_fidelity(source, output, options, workdir)
        return True
    if spec.processor == 'pdf_pptx':
        pdf_to_pptx_fidelity(source, output, options, workdir)
        return True
    if spec.processor == 'pdf_xlsx':
        pdf_to_xlsx_fidelity(source, output, options)
        return True
    if spec.processor == 'pdf_csv':
        pdf_to_csv_fidelity(source, output, options)
        return True
    return False
