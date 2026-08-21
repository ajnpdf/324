from __future__ import annotations
import csv
import html
import io
import json
import os
import re
import shutil
import socket
import ipaddress
import subprocess
import tempfile
import time
import textwrap
import zipfile
import warnings
from dataclasses import dataclass
from email import policy
from email.parser import BytesParser
from pathlib import Path
from typing import Any, Iterable
from functools import lru_cache
from urllib.parse import urljoin, urlparse
from xml.etree import ElementTree as ET
import fitz
import pikepdf
import requests
from bs4 import BeautifulSoup
from docx import Document
from ebooklib import epub
from odf.opendocument import OpenDocumentText
from odf.text import P
from openpyxl import Workbook
from PIL import Image, ImageEnhance, ImageFilter, ImageOps, ImageSequence
from pptx import Presentation
from pptx.util import Inches
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.utils import ImageReader
from reportlab.platypus import Image as RLImage
from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer
try:
    import pillow_heif
    pillow_heif.register_heif_opener()
except Exception:
    pillow_heif = None
try:
    import cairosvg
except Exception:
    cairosvg = None
MAX_IMAGE_PIXELS = max(10000000, int(os.getenv('AJN_MAX_IMAGE_PIXELS', '80000000')))
MAX_IMAGE_FRAMES = max(1, int(os.getenv('AJN_MAX_IMAGE_FRAMES', '120')))
MAX_BATCH_PIXELS = max(50000000, int(os.getenv('AJN_MAX_BATCH_MPIX', '240')) * 1000000)
MAX_PDF_PAGES = max(1, int(os.getenv('AJN_MAX_PDF_PAGES', '300')))
MAX_RENDER_PIXELS = max(50000000, int(os.getenv('AJN_MAX_RENDER_MPIX', '600')) * 1000000)
MAX_OUTPUT_BYTES = max(25, int(os.getenv('AJN_MAX_OUTPUT_MB', '500'))) * 1024 * 1024
Image.MAX_IMAGE_PIXELS = MAX_IMAGE_PIXELS

@dataclass(frozen=True)
class ConversionSpec:
    tool_id: str
    name: str
    category: str
    input_extensions: tuple[str, ...]
    output_extension: str
    output_mime: str
    processor: str
    multi_file: bool = False
    limitation: str | None = None
    external_dependency: str | None = None
SPECS: dict[str, ConversionSpec] = {}

def _add(spec: ConversionSpec) -> None:
    SPECS[spec.tool_id] = spec
for tool_id, name, out_ext, out_mime, processor, limitation in []:
    _add(ConversionSpec(tool_id, name, 'recognized_text', ('.pdf',) if tool_id.startswith('scanned-pdf') else ('.jpg', '.jpeg', '.png', '.webp', '.tif', '.tiff', '.bmp', '.gif', '.heic', '.heif'), out_ext, out_mime, processor, multi_file=tool_id in {*()}, limitation=limitation, external_dependency='' if processor.startswith('recognition_') else None))
for ext, slug, label in [('*', 'image-to-pdf', 'Image to PDF'), ('.jpg', 'jpg-to-pdf', 'JPG to PDF'), ('.jpeg', 'jpeg-to-pdf', 'JPEG to PDF'), ('.png', 'png-to-pdf', 'PNG to PDF'), ('.webp', 'webp-to-pdf', 'WEBP to PDF'), ('.tiff', 'tiff-to-pdf', 'TIFF to PDF'), ('.bmp', 'bmp-to-pdf', 'BMP to PDF'), ('.gif', 'gif-to-pdf', 'GIF to PDF'), ('.svg', 'svg-to-pdf', 'SVG to PDF'), ('.heic', 'heic-to-pdf', 'HEIC to PDF')]:
    inputs = ('.jpg', '.jpeg', '.png', '.webp', '.tif', '.tiff', '.bmp', '.gif', '.svg', '.heic', '.heif') if ext == '*' else ('.tif', '.tiff') if ext == '.tiff' else ('.heic', '.heif') if ext == '.heic' else (ext,)
    _add(ConversionSpec(slug, label, 'image-to-pdf', inputs, '.pdf', 'application/pdf', 'images_to_pdf', multi_file=True, limitation='Animated image formats are converted as individual still frames.' if ext == '.gif' else None))
for slug, label, fmt, limitation in [('pdf-to-image', 'PDF to Image', 'png', None), ('pdf-to-jpg', 'PDF to JPG', 'jpg', None), ('pdf-to-jpeg', 'PDF to JPEG', 'jpeg', None), ('pdf-to-png', 'PDF to PNG', 'png', None), ('pdf-to-webp', 'PDF to WEBP', 'webp', None), ('pdf-to-tiff', 'PDF to TIFF', 'tiff', None), ('pdf-to-bmp', 'PDF to BMP', 'bmp', None), ('pdf-to-gif', 'PDF to GIF', 'gif', 'Multiple pages are returned as an animated GIF.'), ('pdf-to-svg', 'PDF to SVG', 'svg', 'Text and vector fidelity depends on the PDF page content.'), ('pdf-to-avif', 'PDF to AVIF', 'avif', 'AVIF encoding availability depends on the installed Pillow build.'), ('pdf-to-heic', 'PDF to HEIC', 'heic', 'HEIC encoding availability depends on the installed HEIF encoder.')]:
    direct = fmt in {'gif', 'tiff'}
    _add(ConversionSpec(slug, label, 'pdf-to-image', ('.pdf',), f'.{fmt}' if direct else '.zip', 'image/gif' if fmt == 'gif' else 'image/tiff' if fmt == 'tiff' else 'application/zip', f'pdf_to_image:{fmt}', limitation=limitation))
_add(ConversionSpec('pdf-pages-to-zip', 'PDF Pages to ZIP', 'pdf-to-image', ('.pdf',), '.zip', 'application/zip', 'pdf_pages_zip'))
pdf_document_specs = [('pdf-to-word', 'PDF to Word', '.docx', 'application/vnd.openxmlformats-officedocument.wordprocessingml.document', 'pdf_docx', 'Creates an editable text-focused document; complex layout may change.'), ('pdf-to-docx', 'PDF to DOCX', '.docx', 'application/vnd.openxmlformats-officedocument.wordprocessingml.document', 'pdf_docx', 'Creates an editable text-focused document; complex layout may change.'), ('pdf-to-txt', 'PDF to TXT', '.txt', 'text/plain', 'pdf_text', None), ('pdf-to-rtf', 'PDF to RTF', '.rtf', 'application/rtf', 'pdf_rtf', 'Preserves text content rather than exact page design.'), ('pdf-to-odt', 'PDF to ODT', '.odt', 'application/vnd.oasis.opendocument.text', 'pdf_odt', 'Creates an editable text-focused document; complex layout may change.'), ('pdf-to-html', 'PDF to HTML', '.html', 'text/html', 'pdf_html', 'Creates readable HTML from extracted text; advanced PDF positioning may change.'), ('pdf-to-markdown', 'PDF to Markdown', '.md', 'text/markdown', 'pdf_markdown', 'Headings and tables are inferred from extracted text and may require editing.'), ('pdf-to-xml', 'PDF to XML', '.xml', 'application/xml', 'pdf_xml', None), ('pdf-to-json', 'PDF to JSON', '.json', 'application/json', 'pdf_json', None), ('pdf-to-csv', 'PDF to CSV', '.csv', 'text/csv', 'pdf_csv', 'Rows are created from extracted text lines; this is not advanced table recognition.'), ('pdf-to-excel', 'PDF to Excel', '.xlsx', 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 'pdf_xlsx', 'Rows are created from extracted text lines; complex tables may require correction.'), ('pdf-to-xlsx', 'PDF to XLSX', '.xlsx', 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 'pdf_xlsx', 'Rows are created from extracted text lines; complex tables may require correction.'), ('pdf-to-powerpoint', 'PDF to PowerPoint', '.pptx', 'application/vnd.openxmlformats-officedocument.presentationml.presentation', 'pdf_pptx', 'Each PDF page is placed as a slide image to preserve appearance.'), ('pdf-to-pptx', 'PDF to PPTX', '.pptx', 'application/vnd.openxmlformats-officedocument.presentationml.presentation', 'pdf_pptx', 'Each PDF page is placed as a slide image to preserve appearance.'), ('pdf-to-epub', 'PDF to EPUB', '.epub', 'application/epub+zip', 'pdf_epub', 'Creates a reflowable text-focused eBook; complex page design may change.'), ('pdf-to-mobi', 'PDF to MOBI', '.mobi', 'application/x-mobipocket-ebook', 'pdf_ebook_external:mobi', 'Requires Calibre ebook-convert and creates a text-focused eBook.'), ('pdf-to-azw3', 'PDF to AZW3', '.azw3', 'application/vnd.amazon.ebook', 'pdf_ebook_external:azw3', 'Requires Calibre ebook-convert and creates a text-focused eBook.')]
for row in pdf_document_specs:
    _add(ConversionSpec(row[0], row[1], 'pdf-to-document', ('.pdf',), row[2], row[3], row[4], limitation=row[5], external_dependency='calibre' if 'external' in row[4] else None))
for slug, label, exts, processor, limitation, dep in [('word-to-pdf', 'Word to PDF', ('.doc', '.docx'), 'office_to_pdf', 'Complex fonts, macros and advanced layouts may render differently.', 'libreoffice'), ('doc-to-pdf', 'DOC to PDF', ('.doc',), 'office_to_pdf', 'Requires LibreOffice for legacy DOC conversion.', 'libreoffice'), ('docx-to-pdf', 'DOCX to PDF', ('.docx',), 'office_to_pdf', 'Complex fonts and advanced layouts may render differently.', 'libreoffice'), ('txt-to-pdf', 'TXT to PDF', ('.txt',), 'text_document_to_pdf', None, None), ('rtf-to-pdf', 'RTF to PDF', ('.rtf',), 'rtf_to_pdf', 'Rich formatting is simplified for a consistent PDF output.', None), ('odt-to-pdf', 'ODT to PDF', ('.odt',), 'office_to_pdf', 'Requires LibreOffice.', 'libreoffice'), ('ods-to-pdf', 'ODS to PDF', ('.ods',), 'office_to_pdf', 'Requires LibreOffice; print areas and charts may render differently.', 'libreoffice'), ('odp-to-pdf', 'ODP to PDF', ('.odp',), 'office_to_pdf', 'Requires LibreOffice; animations and transitions are not included.', 'libreoffice'), ('html-to-pdf', 'HTML to PDF', ('.html', '.htm'), 'html_file_to_pdf', 'External scripts and complex browser-only CSS are not executed.', None), ('url-to-pdf', 'URL to PDF', tuple(), 'url_to_pdf', 'Creates a readable snapshot; JavaScript-heavy pages may not match the browser exactly.', None), ('markdown-to-pdf', 'Markdown to PDF', ('.md', '.markdown'), 'markdown_to_pdf', None, None), ('xml-to-pdf', 'XML to PDF', ('.xml',), 'structured_text_to_pdf', None, None), ('json-to-pdf', 'JSON to PDF', ('.json',), 'structured_text_to_pdf', None, None), ('csv-to-pdf', 'CSV to PDF', ('.csv',), 'structured_text_to_pdf', None, None), ('excel-to-pdf', 'Excel to PDF', ('.xls', '.xlsx'), 'office_to_pdf', 'Requires LibreOffice; print areas and charts may render differently.', 'libreoffice'), ('xls-to-pdf', 'XLS to PDF', ('.xls',), 'office_to_pdf', 'Requires LibreOffice.', 'libreoffice'), ('xlsx-to-pdf', 'XLSX to PDF', ('.xlsx',), 'office_to_pdf', 'Requires LibreOffice; print areas and charts may render differently.', 'libreoffice'), ('powerpoint-to-pdf', 'PowerPoint to PDF', ('.ppt', '.pptx'), 'office_to_pdf', 'Requires LibreOffice; animations and transitions are not included.', 'libreoffice'), ('ppt-to-pdf', 'PPT to PDF', ('.ppt',), 'office_to_pdf', 'Requires LibreOffice.', 'libreoffice'), ('pptx-to-pdf', 'PPTX to PDF', ('.pptx',), 'office_to_pdf', 'Requires LibreOffice; animations and transitions are not included.', 'libreoffice'), ('epub-to-pdf', 'EPUB to PDF', ('.epub',), 'ebook_to_pdf', 'Creates a printable layout from reflowable eBook content.', None), ('mobi-to-pdf', 'MOBI to PDF', ('.mobi',), 'ebook_external_to_pdf', 'Requires Calibre ebook-convert.', 'calibre'), ('azw3-to-pdf', 'AZW3 to PDF', ('.azw3',), 'ebook_external_to_pdf', 'Requires Calibre ebook-convert.', 'calibre'), ('eml-to-pdf', 'EML to PDF', ('.eml',), 'eml_to_pdf', 'Attachments are listed but are not embedded automatically.', None), ('msg-to-pdf', 'MSG to PDF', ('.msg',), 'msg_to_pdf', 'Attachments are listed but are not embedded automatically.', None), ('xps-to-pdf', 'XPS to PDF', ('.xps',), 'xps_to_pdf', 'Converts standard XPS documents to PDF; complex effects and embedded fonts may render differently.', None)]:
    _add(ConversionSpec(slug, label, 'document-to-pdf', exts, '.pdf', 'application/pdf', processor, limitation=limitation, external_dependency=dep))

def command_path(name: str) -> str | None:
    if name == 'libreoffice' and os.name == 'nt':
        roots = [os.environ.get('ProgramFiles'), os.environ.get('ProgramFiles(x86)'), 'C:\\Program Files', 'C:\\Program Files (x86)']
        candidates: list[Path] = []
        for root in roots:
            if root:
                candidates.append(Path(root) / 'LibreOffice' / 'program' / 'soffice.exe')
        found_exe = shutil.which('soffice.exe')
        if found_exe:
            candidates.insert(0, Path(found_exe))
        for candidate in candidates:
            if candidate.exists() and candidate.name.lower() == 'soffice.exe':
                return str(candidate)
        return None
    candidates = {'libreoffice': ['libreoffice', 'soffice'], 'calibre': ['ebook-convert'], 'mutool': ['mutool', 'gxps', 'gswin64c', 'gs']}.get(name, [name])
    for candidate in candidates:
        found = shutil.which(candidate)
        if found:
            return found
    windows_candidates = {'calibre': [Path('C:\\Program Files\\Calibre2\\ebook-convert.exe')], 'mutool': [Path('C:\\Program Files\\MuPDF\\mutool.exe')]}.get(name, [])
    if name == 'mutool':
        windows_candidates.extend(sorted(Path('C:\\Program Files\\gs').glob('gs*\\bin\\gswin64c.exe'), reverse=True))
        windows_candidates.extend(sorted(Path('C:\\Program Files\\Artifex Software').glob('**\\mutool.exe'), reverse=True))
    for candidate in windows_candidates:
        if candidate.exists():
            return str(candidate)
    return None

def _validate_public_url(raw_url: str) -> str:
    parsed = urlparse(raw_url)
    if parsed.scheme.lower() not in {'http', 'https'}:
        raise ValueError('Only public http and https URLs are accepted.')
    if not parsed.hostname or parsed.username or parsed.password:
        raise ValueError('Enter a valid public URL without embedded credentials.')
    host = parsed.hostname.rstrip('.').lower()
    if host in {'localhost', 'localhost.localdomain'} or host.endswith('.local'):
        raise ValueError('Local and private network URLs are not accepted.')
    try:
        records = socket.getaddrinfo(host, parsed.port or (443 if parsed.scheme == 'https' else 80), type=socket.SOCK_STREAM)
    except socket.gaierror as exc:
        raise ValueError('The URL hostname could not be resolved.') from exc
    for record in records:
        address = ipaddress.ip_address(record[4][0].split('%', 1)[0])
        if not address.is_global:
            raise ValueError('Local, private and reserved network addresses are not accepted.')
    return parsed.geturl()

def _download_public_html(raw_url: str, max_bytes: int=5 * 1024 * 1024) -> tuple[str, str]:
    current = _validate_public_url(raw_url)
    session = requests.Session()
    headers = {'User-Agent': 'AJN-PDF/2.0 (+https://www.ajnpdf.com)'}
    for _ in range(4):
        with session.get(current, timeout=(5, 20), headers=headers, allow_redirects=False, stream=True) as response:
            if response.is_redirect or response.is_permanent_redirect:
                target = response.headers.get('location')
                if not target:
                    raise ValueError('The URL returned an invalid redirect.')
                current = _validate_public_url(urljoin(current, target))
                continue
            response.raise_for_status()
            content_type = response.headers.get('content-type', '').lower()
            if not any((kind in content_type for kind in ('text/html', 'application/xhtml+xml', 'text/plain'))):
                raise ValueError('The URL must return an HTML or text page.')
            declared = response.headers.get('content-length')
            if declared and declared.isdigit() and (int(declared) > max_bytes):
                raise ValueError('The web page is too large to convert safely.')
            chunks: list[bytes] = []
            received = 0
            for chunk in response.iter_content(chunk_size=64 * 1024):
                if not chunk:
                    continue
                received += len(chunk)
                if received > max_bytes:
                    raise ValueError('The web page is too large to convert safely.')
                chunks.append(chunk)
            encoding = response.encoding or 'utf-8'
            return (b''.join(chunks).decode(encoding, errors='replace'), current)
    raise ValueError('The URL redirected too many times.')

def tool_available(spec: ConversionSpec) -> tuple[bool, str | None]:
    dep = spec.external_dependency
    if dep and (not command_path(dep)):
        label = {'libreoffice': 'LibreOffice', 'calibre': 'Calibre ebook-convert', 'mutool': 'a licensed XPS/PostScript conversion engine'}.get(dep, dep)
        return (False, f'{label}is required for this conversion on the processing server.')
    if spec.tool_id == 'svg-to-pdf' and cairosvg is None:
        return (False, 'CairoSVG is required for SVG to PDF conversion.')
    if spec.tool_id in {'heic-to-pdf', 'pdf-to-heic'} and pillow_heif is None:
        return (False, 'A HEIF/HEIC encoder and decoder is required on the processing server.')
    if spec.tool_id == 'pdf-to-avif':
        registered = {extension.lower() for extension in Image.registered_extensions()}
        if '.avif' not in registered:
            return (False, 'AVIF encoding is unavailable in the installed Pillow build.')
    return (True, None)

def list_tools() -> list[dict[str, Any]]:
    result: list[dict[str, Any]] = []
    for spec in SPECS.values():
        available, reason = tool_available(spec)
        result.append({'id': spec.tool_id, 'name': spec.name, 'category': spec.category, 'inputExtensions': list(spec.input_extensions), 'outputExtension': spec.output_extension, 'available': available, 'unavailableReason': reason, 'limitation': spec.limitation, 'multiFile': spec.multi_file, 'processingMode': 'temporary-server'})
    return result

def list_backend_tools() -> list[dict[str, Any]]:
    """Return every backend-dependent public capability used by the frontend build. Protect, Unlock and Repair are intentionally separate API endpoints rather than ConversionSpec entries, so they must be exported here as first-class capabilities. Without these records the dependency-aware frontend build would incorrectly hide those tools even when the backend is healthy."""
    tools = list_tools()
    tools.extend([])
    return tools

def validate_extensions(spec: ConversionSpec, files: list[Path]) -> None:
    if not files and spec.processor != 'url_to_pdf':
        raise ValueError('At least one source file is required.')
    if not spec.multi_file and len(files) > 1:
        raise ValueError('This conversion accepts one file at a time.')
    allowed = set((ext.lower() for ext in spec.input_extensions))
    if allowed:
        for path in files:
            if path.suffix.lower() not in allowed:
                raise ValueError(f"{path.suffix or 'Unknown'}files are not supported by{spec.name}.")
_OLE_SIGNATURE = bytes.fromhex('D0CF11E0A1B11AE1')
_ZIP_BASED_EXPECTATIONS: dict[str, tuple[str, ...]] = {'.docx': ('word/document.xml',), '.xlsx': ('xl/workbook.xml',), '.pptx': ('ppt/presentation.xml',), '.odt': ('content.xml', 'mimetype'), '.ods': ('content.xml', 'mimetype'), '.odp': ('content.xml', 'mimetype'), '.epub': ('META-INF/container.xml',), '.xps': ('[Content_Types].xml',)}
_IMAGE_EXTENSIONS = {'.jpg', '.jpeg', '.png', '.webp', '.tif', '.tiff', '.bmp', '.gif', '.heic', '.heif'}
_TEXT_EXTENSIONS = {'.txt', '.md', '.markdown', '.html', '.htm', '.csv', '.json', '.xml', '.eml', '.rtf', '.svg'}

def _read_probe(path: Path, limit: int=4096) -> bytes:
    with path.open('rb') as handle:
        return handle.read(limit)

def validate_input_file(path: Path, allow_encrypted_pdf: bool=False) -> None:
    """Validate that an uploaded file structurally matches its claimed extension. This intentionally uses conservative, format-aware checks before native/third-party converters receive untrusted data. It is not a malware scanner, but it prevents the common case of renamed or obviously damaged inputs reaching LibreOffice, Pillow, PyMuPDF, Calibre or message parsers."""
    if not path.exists() or path.stat().st_size <= 0:
        raise ValueError('The uploaded file is empty.')
    ext = path.suffix.lower()
    probe = _read_probe(path)
    if ext == '.pdf':
        if not probe.startswith(b'%PDF-'):
            raise ValueError('The uploaded file content is not a valid PDF.')
        try:
            with fitz.open(path) as document:
                if document.page_count < 1:
                    raise ValueError('The PDF contains no pages.')
                if document.page_count > MAX_PDF_PAGES:
                    raise ValueError(f'This PDF has too many pages for one processing job. Maximum:{MAX_PDF_PAGES}. Split the PDF and try again.')
                if document.needs_pass:
                    if allow_encrypted_pdf:
                        return
                    raise ValueError('The PDF is password protected. Unlock it first, then try this tool again.')
                document.load_page(0)
        except ValueError:
            raise
        except Exception as exc:
            raise ValueError('The uploaded PDF is damaged or unreadable.') from exc
        return
    if ext in _IMAGE_EXTENSIONS:
        try:
            with Image.open(path) as image:
                if image.width * image.height > MAX_IMAGE_PIXELS:
                    raise ValueError('The image dimensions are too large to process safely.')
                image.verify()
        except Exception as exc:
            raise ValueError('The uploaded image is damaged or does not match its file type.') from exc
        return
    if ext == '.svg':
        try:
            raw_svg = path.read_text(encoding='utf-8', errors='strict')
            root = ET.fromstring(raw_svg)
            if not root.tag.lower().endswith('svg'):
                raise ValueError
            for element in root.iter():
                for attr_name, attr_value in element.attrib.items():
                    if attr_name.lower().endswith('href'):
                        value = str(attr_value).strip()
                        if value and (not value.startswith('#')) and (not value.lower().startswith('data:')):
                            raise ValueError('External SVG resources are not allowed.')
            for match in re.findall('url\\(([^)]+)\\)', raw_svg, flags=re.I):
                value = match.strip().strip('"\'')
                if value and (not value.startswith('#')) and (not value.lower().startswith('data:')):
                    raise ValueError('External SVG resources are not allowed.')
            if re.search('@import\\s+(?:url\\()?\\s*[\\"\']?(?!data:|#)', raw_svg, flags=re.I):
                raise ValueError('External SVG resources are not allowed.')
        except ValueError:
            raise
        except Exception as exc:
            raise ValueError('The uploaded SVG is not valid XML/SVG content.') from exc
        return
    if ext in _ZIP_BASED_EXPECTATIONS:
        if not zipfile.is_zipfile(path):
            raise ValueError(f"The uploaded{ext.upper().lstrip('.')}file is damaged or has the wrong file type.")
        try:
            with zipfile.ZipFile(path) as archive:
                names = set(archive.namelist())
                missing = [name for name in _ZIP_BASED_EXPECTATIONS[ext] if name not in names]
                if missing:
                    raise ValueError(f"The uploaded{ext.upper().lstrip('.')}file is missing required document data.")
                if len(names) > 20000:
                    raise ValueError('The uploaded document contains too many archive entries to process safely.')
                declared = sum((max(0, item.file_size) for item in archive.infolist()))
                if declared > 512 * 1024 * 1024:
                    raise ValueError('The expanded document is too large to process safely.')
        except ValueError:
            raise
        except Exception as exc:
            raise ValueError(f"The uploaded{ext.upper().lstrip('.')}document container is damaged.") from exc
        return
    if ext in {'.doc', '.xls', '.ppt', '.msg'}:
        if not probe.startswith(_OLE_SIGNATURE):
            raise ValueError(f"The uploaded{ext.upper().lstrip('.')}file does not match the expected document format.")
        return
    if ext == '.rtf':
        if not probe.lstrip().lower().startswith(b'{\\rtf'):
            raise ValueError('The uploaded RTF file is not valid RTF content.')
        return
    if ext == '.json':
        try:
            json.loads(path.read_text(encoding='utf-8-sig'))
        except Exception as exc:
            raise ValueError('The uploaded JSON file is not valid JSON.') from exc
        return
    if ext == '.xml':
        try:
            ET.parse(path)
        except Exception as exc:
            raise ValueError('The uploaded XML file is not valid XML.') from exc
        return
    if ext == '.csv':
        try:
            raw = path.read_text(encoding='utf-8-sig', errors='strict')
            rows = list(csv.reader(io.StringIO(raw)))
            if not rows:
                raise ValueError
        except Exception as exc:
            raise ValueError('The uploaded CSV file could not be read safely.') from exc
        return
    if ext == '.eml':
        try:
            message = BytesParser(policy=policy.default).parsebytes(path.read_bytes())
            if not any((message.get('from'), message.get('to'), message.get('subject'), message.get_payload())):
                raise ValueError
        except Exception as exc:
            raise ValueError('The uploaded EML message is not readable.') from exc
        return
    if ext in {'.mobi', '.azw3'}:
        data = path.read_bytes()[:96]
        if len(data) < 68 or data[60:68] not in {b'BOOKMOBI', b'TEXtREAd'}:
            raise ValueError(f"The uploaded{ext.upper().lstrip('.')}eBook does not match the expected format.")
        return
    if ext in _TEXT_EXTENSIONS or ext in {'.html', '.htm'}:
        try:
            text = path.read_text(encoding='utf-8-sig', errors='strict')
            if '\x00' in text:
                raise ValueError
        except Exception as exc:
            raise ValueError('The uploaded text document could not be decoded safely.') from exc

def validate_input_files(spec: ConversionSpec, files: list[Path]) -> None:
    validate_extensions(spec, files)
    for path in files:
        validate_input_file(path)

def validate_output_file(path: Path, expected_extension: str, allow_encrypted_pdf: bool=False) -> None:
    if not path.exists() or path.stat().st_size < 1:
        raise RuntimeError('The converter did not create an output file.')
    if path.stat().st_size > MAX_OUTPUT_BYTES:
        raise RuntimeError('The generated output is too large to return safely. Split the source or reduce output quality.')
    if path.suffix.lower() != expected_extension.lower():
        raise RuntimeError('The converter returned an unexpected output type.')
    ext = path.suffix.lower()
    probe = _read_probe(path, 32)
    if ext == '.pdf':
        if not probe.startswith(b'%PDF-'):
            raise RuntimeError('The generated PDF is invalid.')
        try:
            with fitz.open(path) as document:
                if document.page_count < 1:
                    raise RuntimeError('The generated PDF contains no pages.')
                if document.needs_pass:
                    if not allow_encrypted_pdf:
                        raise RuntimeError('The generated PDF is unexpectedly encrypted.')
                else:
                    document.load_page(0)
        except RuntimeError:
            raise
        except Exception as exc:
            raise RuntimeError('The generated PDF is damaged.') from exc
    elif ext in {'.zip', '.docx', '.xlsx', '.pptx', '.odt', '.ods', '.odp', '.epub'}:
        if not zipfile.is_zipfile(path):
            raise RuntimeError(f"The generated{ext.upper().lstrip('.')}file is not a valid container.")
        with zipfile.ZipFile(path) as archive:
            names = set(archive.namelist())
            if not names:
                raise RuntimeError('The generated archive is empty.')
            if ext in _ZIP_BASED_EXPECTATIONS:
                missing = [name for name in _ZIP_BASED_EXPECTATIONS[ext] if name not in names]
                if missing:
                    raise RuntimeError(f"The generated{ext.upper().lstrip('.')}file is missing required document data.")
    elif ext in _IMAGE_EXTENSIONS or ext == '.avif':
        try:
            with Image.open(path) as image:
                if image.width * image.height > MAX_IMAGE_PIXELS:
                    raise RuntimeError('The generated image dimensions exceed the safe output limit.')
                image.verify()
        except RuntimeError:
            raise
        except Exception as exc:
            raise RuntimeError('The generated image is invalid.') from exc
    elif ext == '.svg':
        try:
            root = ET.parse(path).getroot()
            if not root.tag.lower().endswith('svg'):
                raise ValueError
        except Exception as exc:
            raise RuntimeError('The generated SVG is invalid.') from exc
    elif ext == '.json':
        try:
            json.loads(path.read_text(encoding='utf-8-sig'))
        except Exception as exc:
            raise RuntimeError('The generated JSON is invalid.') from exc
    elif ext == '.xml':
        try:
            ET.parse(path)
        except Exception as exc:
            raise RuntimeError('The generated XML is invalid.') from exc
    elif ext == '.csv':
        try:
            with path.open('r', encoding='utf-8-sig', errors='strict', newline='') as handle:
                if next(csv.reader(handle), None) is None:
                    raise ValueError
        except Exception as exc:
            raise RuntimeError('The generated CSV is invalid.') from exc
    elif ext == '.rtf':
        if not probe.lstrip().lower().startswith(b'{\\rtf'):
            raise RuntimeError('The generated RTF is invalid.')
    elif ext in {'.mobi', '.azw3'}:
        data = path.read_bytes()[:96]
        if len(data) < 68 or data[60:68] not in {b'BOOKMOBI', b'TEXtREAd'}:
            raise RuntimeError(f"The generated{ext.upper().lstrip('.')}eBook is invalid.")
    elif ext in _TEXT_EXTENSIONS or ext in {'.html', '.htm'}:
        try:
            text = path.read_text(encoding='utf-8-sig', errors='strict')
            if '\x00' in text:
                raise ValueError
        except Exception as exc:
            raise RuntimeError('The generated text output could not be decoded safely.') from exc

def _open_image(path: Path) -> Image.Image:
    if path.suffix.lower() == '.svg':
        if cairosvg is None:
            raise RuntimeError('SVG conversion requires CairoSVG on the server.')
        png_bytes = cairosvg.svg2png(bytestring=path.read_bytes())
        image = Image.open(io.BytesIO(png_bytes))
    else:
        image = Image.open(path)
    image.load()
    return image

def _rgb(image: Image.Image, background: str='white') -> Image.Image:
    if image.mode in {'RGBA', 'LA'}:
        base = Image.new('RGB', image.size, background)
        alpha = image.getchannel('A') if 'A' in image.getbands() else image.getchannel('L')
        base.paste(image.convert('RGB'), mask=alpha)
        return base
    return image.convert('RGB')

def _iter_frames(path: Path) -> Iterable[Image.Image]:
    image = _open_image(path)
    frame_count = 0
    total_pixels = 0
    try:
        for frame in ImageSequence.Iterator(image):
            frame_count += 1
            if frame_count > MAX_IMAGE_FRAMES:
                raise ValueError(f'This image contains too many frames. Maximum:{MAX_IMAGE_FRAMES}.')
            frame_pixels = max(1, frame.width) * max(1, frame.height)
            if frame_pixels > MAX_IMAGE_PIXELS:
                raise ValueError('An image frame is too large to process safely.')
            total_pixels += frame_pixels
            if total_pixels > MAX_BATCH_PIXELS:
                raise ValueError('The image frames are too large for one processing job. Use fewer or smaller images.')
            yield _rgb(frame.copy())
    finally:
        image.close()

def _preprocess_scan(image: Image.Image, grayscale: bool=True) -> Image.Image:
    image = ImageOps.exif_transpose(image)
    if grayscale:
        image = ImageOps.grayscale(image)
        image = ImageOps.autocontrast(image, cutoff=1)
        image = ImageEnhance.Contrast(image).enhance(1.25)
        image = image.filter(ImageFilter.SHARPEN)
        return image.convert('RGB')
    return _rgb(ImageOps.autocontrast(image.convert('RGB'), cutoff=1))

def _images_to_pdf(files: list[Path], output: Path, options: dict[str, Any], preprocess: bool=False) -> None:
    images: list[Image.Image] = []
    total_pixels = 0
    total_frames = 0
    try:
        for file in files:
            for frame in _iter_frames(file):
                total_frames += 1
                if total_frames > MAX_IMAGE_FRAMES:
                    frame.close()
                    raise ValueError(f'Too many image pages for one PDF. Maximum:{MAX_IMAGE_FRAMES}.')
                total_pixels += max(1, frame.width) * max(1, frame.height)
                if total_pixels > MAX_BATCH_PIXELS:
                    frame.close()
                    raise ValueError('The selected images are too large for one PDF job. Use fewer or smaller images.')
                processed = _preprocess_scan(frame, bool(options.get('grayscale', preprocess))) if preprocess else _rgb(ImageOps.exif_transpose(frame))
                if processed is not frame:
                    frame.close()
                images.append(processed)
        if not images:
            raise ValueError('No readable image frames were found.')
        quality = max(40, min(100, int(options.get('quality', 90))))
        first, rest = (images[0], images[1:])
        first.save(output, 'PDF', save_all=True, append_images=rest, resolution=float(options.get('dpi', 150)), quality=quality)
    finally:
        for image in images:
            image.close()

def _pdf_render_workload(document: fitz.Document, dpi: int) -> int:
    if document.page_count < 1:
        raise ValueError('The PDF has no pages.')
    if document.page_count > MAX_PDF_PAGES:
        raise ValueError(f'This PDF has too many pages for one processing job. Maximum:{MAX_PDF_PAGES}. Split the PDF and try again.')
    total_pixels = 0
    scale = dpi / 72
    for page in document:
        rect = page.rect
        total_pixels += max(1, int(rect.width * scale)) * max(1, int(rect.height * scale))
        if total_pixels > MAX_RENDER_PIXELS:
            raise ValueError('The selected page count and resolution are too large for one job. Lower the resolution or split the PDF.')
    return total_pixels

def _iter_pdf_rendered_pages(pdf_path: Path, dpi: int=150) -> Iterable[tuple[int, Image.Image]]:
    dpi = max(72, min(600, int(dpi)))
    document = fitz.open(pdf_path)
    try:
        _pdf_render_workload(document, dpi)
        matrix = fitz.Matrix(dpi / 72, dpi / 72)
        for index, page in enumerate(document, start=1):
            pix = page.get_pixmap(matrix=matrix, alpha=False)
            image = Image.frombytes('RGB', [pix.width, pix.height], pix.samples)
            yield (index, image)
    finally:
        document.close()

def _render_pdf_pages(pdf_path: Path, dpi: int=150) -> list[Image.Image]:
    images: list[Image.Image] = []
    total_pixels = 0
    try:
        for _, image in _iter_pdf_rendered_pages(pdf_path, dpi):
            if len(images) >= MAX_IMAGE_FRAMES:
                image.close()
                raise ValueError(f'This output requires too many in-memory frames. Maximum:{MAX_IMAGE_FRAMES}.')
            total_pixels += max(1, image.width) * max(1, image.height)
            if total_pixels > MAX_BATCH_PIXELS:
                image.close()
                raise ValueError('This GIF/TIFF job is too large to keep safely in memory. Lower the resolution or split the PDF.')
            images.append(image)
        return images
    except Exception:
        for image in images:
            image.close()
        raise

def _zip_files(paths: list[Path], output: Path) -> None:
    with zipfile.ZipFile(output, 'w', zipfile.ZIP_DEFLATED) as archive:
        for path in paths:
            archive.write(path, arcname=path.name)

def _pdf_to_images(source: Path, output: Path, fmt: str, options: dict[str, Any], workdir: Path) -> None:
    dpi = max(72, min(600, int(options.get('dpi', 150))))
    quality = max(35, min(100, int(options.get('quality', 90))))
    generated: list[Path] = []
    if fmt == 'svg':
        document = fitz.open(source)
        try:
            if document.page_count < 1:
                raise ValueError('The PDF has no pages.')
            if document.page_count > MAX_PDF_PAGES:
                raise ValueError(f'This PDF has too many pages for one processing job. Maximum:{MAX_PDF_PAGES}. Split the PDF and try again.')
            for index, page in enumerate(document, start=1):
                path = workdir / f'page-{index:03d}.svg'
                path.write_text(page.get_svg_image(), encoding='utf-8')
                generated.append(path)
        finally:
            document.close()
        _zip_files(generated, output)
        return
    if fmt in {'gif', 'tiff'}:
        images = _render_pdf_pages(source, dpi)
        if not images:
            raise ValueError('The PDF has no pages.')
        try:
            if fmt == 'gif':
                images[0].save(output, 'GIF', save_all=True, append_images=images[1:], duration=int(options.get('frame_duration_ms', 700)), loop=0)
            else:
                images[0].save(output, 'TIFF', save_all=True, append_images=images[1:], compression='tiff_deflate')
        finally:
            for image in images:
                image.close()
        return
    for index, image in _iter_pdf_rendered_pages(source, dpi):
        normalized_fmt = 'JPEG' if fmt in {'jpg', 'jpeg'} else 'HEIF' if fmt == 'heic' else fmt.upper()
        suffix = 'jpg' if fmt in {'jpg', 'jpeg'} else fmt
        path = workdir / f'page-{index:03d}.{suffix}'
        kwargs: dict[str, Any] = {}
        if normalized_fmt in {'JPEG', 'WEBP', 'AVIF', 'HEIF'}:
            kwargs['quality'] = quality
        try:
            image.save(path, normalized_fmt, **kwargs)
        except Exception as exc:
            raise RuntimeError(f'{fmt.upper()}encoding is not available on this server.') from exc
        finally:
            image.close()
        generated.append(path)
    _zip_files(generated, output)

def _extract_pdf_pages(source: Path) -> list[dict[str, Any]]:
    doc = fitz.open(source)
    if doc.page_count > MAX_PDF_PAGES:
        doc.close()
        raise ValueError(f'This PDF has too many pages for one processing job. Maximum:{MAX_PDF_PAGES}. Split the PDF and try again.')
    pages: list[dict[str, Any]] = []
    for index, page in enumerate(doc, start=1):
        text = page.get_text('text').strip()
        pages.append({'page': index, 'text': text, 'lines': [line for line in text.splitlines() if line.strip()]})
    doc.close()
    return pages

def _write_docx(pages: list[dict[str, Any]], output: Path, title: str) -> None:
    doc = Document()
    doc.core_properties.title = title
    doc.add_heading(title, level=0)
    for page in pages:
        doc.add_heading(f"Page{page['page']}", level=1)
        text = page['text'] or '[No text detected on this page]'
        for paragraph in re.split('\\n\\s*\\n', text):
            doc.add_paragraph(paragraph.strip())
    doc.save(output)

def _write_pdf_text(pages: list[dict[str, Any]], output: Path, title: str) -> None:
    styles = getSampleStyleSheet()
    story: list[Any] = [Paragraph(html.escape(title), styles['Title']), Spacer(1, 12)]
    for page in pages:
        story.append(Paragraph(f"Page{page['page']}", styles['Heading2']))
        text = page['text'] or 'No text detected.'
        for paragraph in re.split('\\n\\s*\\n', text):
            safe = html.escape(paragraph.strip()).replace('', '<br/>')
            if safe:
                story.append(Paragraph(safe, styles['BodyText']))
                story.append(Spacer(1, 8))
        story.append(PageBreak())
    SimpleDocTemplate(str(output), pagesize=A4, rightMargin=42, leftMargin=42, topMargin=42, bottomMargin=42).build(story)

def _pdf_to_pptx(source: Path, output: Path, options: dict[str, Any], workdir: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    slide_count = 0
    for index, image in _iter_pdf_rendered_pages(source, int(options.get('dpi', 150))):
        path = workdir / f'slide-{index:03d}.png'
        try:
            image.save(path, 'PNG')
        finally:
            image.close()
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(path), 0, 0, width=prs.slide_width, height=prs.slide_height)
        slide_count += 1
    if slide_count < 1:
        raise ValueError('The PDF has no pages.')
    prs.save(output)

def _pdf_to_epub(pages: list[dict[str, Any]], output: Path, title: str) -> None:
    book = epub.EpubBook()
    book.set_identifier(f'ajn-{abs(hash(title))}')
    book.set_title(title)
    book.set_language('en')
    chapters = []
    for page in pages:
        chapter = epub.EpubHtml(title=f"Page{page['page']}", file_name=f"page-{page['page']}.xhtml", lang='en')
        chapter.content = f"<h1>Page{page['page']}</h1><pre>{html.escape(page['text'])}</pre>"
        book.add_item(chapter)
        chapters.append(chapter)
    book.toc = tuple(chapters)
    book.spine = ['nav', *chapters]
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    epub.write_epub(str(output), book)

def _terminate_process_tree(process: subprocess.Popen[bytes]) -> None:
    if process.poll() is not None:
        return
    if os.name == 'nt':
        try:
            subprocess.run(['taskkill', '/PID', str(process.pid), '/T', '/F'], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, timeout=15, check=False)
            return
        except Exception:
            pass
    try:
        process.kill()
    except Exception:
        pass

def _run(command: list[str], timeout: int=180) -> None:
    creationflags = 0
    if os.name == 'nt':
        creationflags = getattr(subprocess, 'CREATE_NEW_PROCESS_GROUP', 0)
    process = subprocess.Popen(command, stdout=subprocess.PIPE, stderr=subprocess.PIPE, creationflags=creationflags)
    try:
        stdout, stderr = process.communicate(timeout=max(1, int(timeout)))
    except subprocess.TimeoutExpired as exc:
        _terminate_process_tree(process)
        try:
            stdout, stderr = process.communicate(timeout=10)
        except Exception:
            stdout, stderr = (b'', b'')
        detail = (stderr or stdout).decode('utf-8', errors='ignore')[-800:]
        message = f'Command timed out after{int(timeout)}seconds.'
        if detail.strip():
            message += f'{detail.strip()}'
        raise RuntimeError(message) from exc
    if process.returncode != 0:
        detail = stderr.decode('utf-8', errors='ignore')[-800:] or stdout.decode('utf-8', errors='ignore')[-800:]
        raise RuntimeError(detail.strip() or f'Command failed with exit code{process.returncode}.')

def _office_to_pdf(source: Path, output: Path, workdir: Path) -> None:
    executable = command_path('libreoffice')
    if not executable:
        raise RuntimeError('LibreOffice is not installed on the processing server.')
    last_error: Exception | None = None
    for attempt in range(1, 3):
        profile_dir = workdir / f'libreoffice-profile-{attempt}'
        if profile_dir.exists():
            shutil.rmtree(profile_dir, ignore_errors=True)
        profile_dir.mkdir(parents=True, exist_ok=True)
        profile_uri = profile_dir.resolve().as_uri()
        generated = workdir / f'{source.stem}.pdf'
        generated.unlink(missing_ok=True)
        try:
            _run([executable, '--headless', '--nologo', '--nodefault', '--nolockcheck', '--norestore', f'-env:UserInstallation={profile_uri}', '--convert-to', 'pdf', '--outdir', str(workdir), str(source)], timeout=150)
            previous_size = -1
            stable_samples = 0
            for _ in range(60):
                if generated.exists() and generated.stat().st_size > 32:
                    current_size = generated.stat().st_size
                    if current_size == previous_size:
                        stable_samples += 1
                        if stable_samples >= 2:
                            break
                    else:
                        stable_samples = 0
                    previous_size = current_size
                time.sleep(0.25)
            if not generated.exists() or generated.stat().st_size <= 32:
                candidates = [p for p in workdir.glob('*.pdf') if p != output and p.stat().st_size > 32]
                if not candidates:
                    raise RuntimeError('LibreOffice did not create a PDF output.')
                generated = max(candidates, key=lambda item: item.stat().st_mtime)
            try:
                with fitz.open(generated) as office_pdf:
                    if office_pdf.page_count < 1:
                        raise RuntimeError('LibreOffice created a PDF with no pages.')
                    office_pdf.load_page(0)
            except RuntimeError:
                raise
            except Exception as exc:
                raise RuntimeError('LibreOffice created an invalid or incomplete PDF output.') from exc
            shutil.move(str(generated), str(output))
            return
        except Exception as exc:
            last_error = exc
            if attempt == 1:
                time.sleep(1)
                continue
            break
        finally:
            shutil.rmtree(profile_dir, ignore_errors=True)
    raise RuntimeError(f'LibreOffice conversion failed after retry:{last_error}')

def _ebook_external(source: Path, output: Path) -> None:
    executable = command_path('calibre')
    if not executable:
        raise RuntimeError('Calibre ebook-convert is not installed on the processing server.')
    _run([executable, str(source), str(output)], timeout=300)

def _xps_to_pdf(source: Path, output: Path) -> None:
    executable = command_path('mutool')
    if not executable:
        raise RuntimeError('MuPDF or Ghostscript is not installed on the processing server.')
    name = Path(executable).name.lower()
    if 'mutool' in name:
        _run([executable, 'convert', '-o', str(output), str(source)], timeout=240)
    elif 'gxps' in name:
        _run([executable, '-sDEVICE=pdfwrite', f'-sOutputFile={output}', str(source)], timeout=240)
    else:
        _run([executable, '-dBATCH', '-dNOPAUSE', '-sDEVICE=pdfwrite', f'-sOutputFile={output}', str(source)], timeout=240)

def _html_to_text(raw: str) -> str:
    soup = BeautifulSoup(raw, 'html.parser')
    for tag in soup(['script', 'style', 'noscript']):
        tag.decompose()
    return soup.get_text('', strip=True)

def _read_structured_text(path: Path) -> str:
    raw = path.read_text(encoding='utf-8', errors='replace')
    if path.suffix.lower() == '.json':
        try:
            return json.dumps(json.loads(raw), indent=2, ensure_ascii=False)
        except Exception:
            return raw
    if path.suffix.lower() == '.xml':
        return BeautifulSoup(raw, 'xml').prettify()
    if path.suffix.lower() == '.csv':
        rows = list(csv.reader(io.StringIO(raw)))
        return ''.join(('|'.join(row) for row in rows))
    return raw

def _validate_generated_pdf_bytes(pdf_bytes: bytes) -> None:
    if len(pdf_bytes) < 64 or not pdf_bytes.startswith(b'%PDF-'):
        raise RuntimeError('returned invalid searchable PDF data.')
    try:
        with fitz.open(stream=pdf_bytes, filetype='pdf') as generated:
            if generated.page_count != 1:
                raise RuntimeError('returned an invalid searchable PDF page count.')
            generated.load_page(0)
    except RuntimeError:
        raise
    except Exception as exc:
        raise RuntimeError('returned a damaged searchable PDF page.') from exc

def _pdf_pages_zip(source: Path, output: Path, workdir: Path) -> None:
    doc = fitz.open(source)
    paths: list[Path] = []
    try:
        if doc.page_count < 1:
            raise ValueError('The PDF has no pages.')
        if doc.page_count > MAX_PDF_PAGES:
            raise ValueError(f'This PDF has too many pages for one processing job. Maximum:{MAX_PDF_PAGES}. Split the PDF and try again.')
        for index in range(doc.page_count):
            part = fitz.open()
            try:
                part.insert_pdf(doc, from_page=index, to_page=index)
                path = workdir / f'page-{index + 1:03d}.pdf'
                part.save(path)
                paths.append(path)
            finally:
                part.close()
    finally:
        doc.close()
    _zip_files(paths, output)

def convert(spec: ConversionSpec, files: list[Path], output: Path, options: dict[str, Any], workdir: Path, source_url: str | None=None) -> None:
    if spec.processor == 'xps_to_pdf':
        if not files:
            raise ValueError('XPS to PDF requires one XPS file.')
        try:
            import pymupdf as _ajn_pymupdf
            _ajn_xps_doc = _ajn_pymupdf.open(str(files[0]))
            try:
                _ajn_pdf_bytes = _ajn_xps_doc.convert_to_pdf()
            finally:
                _ajn_xps_doc.close()
            if not _ajn_pdf_bytes.startswith(b'%PDF-'):
                raise ValueError('The XPS document could not be converted to a valid PDF.')
            with open(output, 'wb') as _ajn_out:
                _ajn_out.write(_ajn_pdf_bytes)
            return
        except ValueError:
            raise
        except Exception as _ajn_exc:
            raise ValueError('The XPS document could not be converted.') from _ajn_exc
    validate_input_files(spec, files)
    available, reason = tool_available(spec)
    if not available:
        raise RuntimeError(reason or 'This conversion is unavailable on the processing server.')
    processor = spec.processor
    source = files[0] if files else None
    if processor == 'images_to_pdf':
        _images_to_pdf(files, output, options)
    elif processor == 'scan_images_pdf':
        _images_to_pdf(files, output, options, preprocess=True)
    elif processor.startswith('pdf_to_image:'):
        _pdf_to_images(source, output, processor.split(':', 1)[1], options, workdir)
    elif processor == 'pdf_pages_zip':
        _pdf_pages_zip(source, output, workdir)
    elif processor in {'pdf_text', 'pdf_docx', 'pdf_rtf', 'pdf_odt', 'pdf_html', 'pdf_markdown', 'pdf_xml', 'pdf_json', 'pdf_csv', 'pdf_xlsx', 'pdf_pptx', 'pdf_epub'}:
        pages = _extract_pdf_pages(source)
        title = source.stem
        if processor == 'pdf_text':
            output.write_text(''.join((f"--- Page{p['page']}---{p['text']}" for p in pages)), encoding='utf-8')
        elif processor == 'pdf_docx':
            _write_docx(pages, output, title)
        elif processor == 'pdf_rtf':
            body = '\\par'.join((re.sub('([\\\\{}])', '\\\\\\1', p['text']) for p in pages))
            output.write_text('{\\rtf1\\ansi' + body + '}', encoding='utf-8')
        elif processor == 'pdf_odt':
            doc = OpenDocumentText()
            for page in pages:
                doc.text.addElement(P(text=f"Page{page['page']}"))
                for line in page['lines']:
                    doc.text.addElement(P(text=line))
            doc.save(str(output))
        elif processor == 'pdf_html':
            output.write_text("<!doctype html><meta charset='utf-8'><title>" + html.escape(title) + '</title>' + ''.join((f"<section><h2>Page{p['page']}</h2><pre>{html.escape(p['text'])}</pre></section>" for p in pages)), encoding='utf-8')
        elif processor == 'pdf_markdown':
            output.write_text(''.join((f"## Page{p['page']}{p['text']}" for p in pages)), encoding='utf-8')
        elif processor == 'pdf_xml':
            xml = ['<?xml version="1.0" encoding="UTF-8"?>', '<document>']
            for p in pages:
                xml.append(f'''<page number="{p['page']}"><![CDATA[{p['text']}]]></page>''')
            xml.append('</document>')
            output.write_text(''.join(xml), encoding='utf-8')
        elif processor == 'pdf_json':
            output.write_text(json.dumps({'title': title, 'pages': pages}, ensure_ascii=False, indent=2), encoding='utf-8')
        elif processor == 'pdf_csv':
            with output.open('w', newline='', encoding='utf-8-sig') as handle:
                writer = csv.writer(handle)
                writer.writerow(['page', 'line', 'text'])
                for p in pages:
                    for line_number, line in enumerate(p['lines'], start=1):
                        writer.writerow([p['page'], line_number, line])
        elif processor == 'pdf_xlsx':
            workbook = Workbook()
            sheet = workbook.active
            sheet.title = 'PDF text'
            sheet.append(['Page', 'Line', 'Text'])
            for p in pages:
                for line_number, line in enumerate(p['lines'], start=1):
                    sheet.append([p['page'], line_number, line])
            workbook.save(output)
        elif processor == 'pdf_pptx':
            _pdf_to_pptx(source, output, options, workdir)
        elif processor == 'pdf_epub':
            _pdf_to_epub(pages, output, title)
    elif processor.startswith('pdf_ebook_external:'):
        pages = _extract_pdf_pages(source)
        epub_path = workdir / 'intermediate.epub'
        _pdf_to_epub(pages, epub_path, source.stem)
        _ebook_external(epub_path, output)
    elif processor == 'office_to_pdf':
        _office_to_pdf(source, output, workdir)
    elif processor in {'text_document_to_pdf', 'structured_text_to_pdf'}:
        text = _read_structured_text(source)
        _write_pdf_text([{'page': 1, 'text': text, 'lines': text.splitlines()}], output, source.stem)
    elif processor == 'rtf_to_pdf':
        from striprtf.striprtf import rtf_to_text
        text = rtf_to_text(source.read_text(encoding='utf-8', errors='replace'))
        _write_pdf_text([{'page': 1, 'text': text, 'lines': text.splitlines()}], output, source.stem)
    elif processor == 'markdown_to_pdf':
        import markdown
        raw = source.read_text(encoding='utf-8', errors='replace')
        text = _html_to_text(markdown.markdown(raw))
        _write_pdf_text([{'page': 1, 'text': text, 'lines': text.splitlines()}], output, source.stem)
    elif processor == 'html_file_to_pdf':
        raw = source.read_text(encoding='utf-8', errors='replace')
        text = _html_to_text(raw)
        _write_pdf_text([{'page': 1, 'text': text, 'lines': text.splitlines()}], output, source.stem)
    elif processor == 'url_to_pdf':
        if not source_url or not re.match('^https?://', source_url, flags=re.I):
            raise ValueError('Enter a valid http or https URL.')
        html_body, final_url = _download_public_html(source_url)
        text = _html_to_text(html_body)
        if not text.strip():
            raise ValueError('The web page did not contain readable text.')
        _write_pdf_text([{'page': 1, 'text': text, 'lines': text.splitlines()}], output, final_url)
    elif processor == 'ebook_to_pdf':
        with warnings.catch_warnings():
            warnings.filterwarnings('ignore', message='In the future version we will turn default option ignore_ncx to True\\.', category=UserWarning, module='ebooklib\\.epub')
            warnings.filterwarnings('ignore', message='This search incorrectly ignores the root element, and will be fixed in a future version\\..*', category=FutureWarning, module='ebooklib\\.epub')
            book = epub.read_epub(str(source), options={'ignore_ncx': True})
        sections = []
        for item in book.get_items():
            if item.get_type() == 9:
                sections.append(_html_to_text(item.get_content().decode('utf-8', errors='replace')))
        text = ''.join((section for section in sections if section))
        _write_pdf_text([{'page': 1, 'text': text, 'lines': text.splitlines()}], output, source.stem)
    elif processor == 'ebook_external_to_pdf':
        _ebook_external(source, output)
    elif processor == 'eml_to_pdf':
        message = BytesParser(policy=policy.default).parsebytes(source.read_bytes())
        body = message.get_body(preferencelist=('plain', 'html'))
        body_text = ''
        if body:
            content = body.get_content()
            body_text = _html_to_text(content) if body.get_content_type() == 'text/html' else str(content)
        attachments = [part.get_filename() for part in message.iter_attachments() if part.get_filename()]
        text = f"From:{message.get('from', '')}To:{message.get('to', '')}Subject:{message.get('subject', '')}Date:{message.get('date', '')}{body_text}"
        if attachments:
            text += 'Attachments:' + ''.join((f'-{name}' for name in attachments))
        _write_pdf_text([{'page': 1, 'text': text, 'lines': text.splitlines()}], output, source.stem)
    elif processor == 'msg_to_pdf':
        import extract_msg
        msg = extract_msg.Message(str(source))
        text = f"From:{msg.sender or ''}To:{msg.to or ''}Subject:{msg.subject or ''}Date:{msg.date or ''}{msg.body or ''}"
        attachments = [getattr(a, 'longFilename', None) or getattr(a, 'shortFilename', None) for a in msg.attachments]
        attachments = [name for name in attachments if name]
        if attachments:
            text += 'Attachments:' + ''.join((f'-{name}' for name in attachments))
        _write_pdf_text([{'page': 1, 'text': text, 'lines': text.splitlines()}], output, source.stem)
        msg.close()
    elif processor == 'xps_to_pdf':
        _xps_to_pdf(source, output)
    else:
        raise RuntimeError(f'No conversion processor is registered for{spec.tool_id}.')
    validate_output_file(output, spec.output_extension)
