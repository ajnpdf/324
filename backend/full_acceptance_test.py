from __future__ import annotations
import json
import shutil
import subprocess
import tempfile
from pathlib import Path
import fitz
from docx import Document
from ebooklib import epub
from openpyxl import Workbook
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from reportlab.pdfgen import canvas
from app.conversion_engine import SPECS, command_path, convert, tool_available

def _libreoffice_fixture(source: Path, target_extension: str, filter_name: str, root: Path) -> Path | None:
    executable = command_path('libreoffice')
    if not executable:
        return None
    target = source.with_suffix(target_extension)
    profile = root / f"fixture-profile-{target_extension.lstrip('.')}"
    shutil.rmtree(profile, ignore_errors=True)
    profile.mkdir(parents=True, exist_ok=True)
    try:
        completed = subprocess.run([executable, '--headless', '--nologo', '--nodefault', '--nolockcheck', '--norestore', f'-env:UserInstallation={profile.resolve().as_uri()}', '--convert-to', filter_name, '--outdir', str(root), str(source)], stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=120, check=False)
        if completed.returncode == 0 and target.exists() and (target.stat().st_size > 32):
            return target
        return None
    finally:
        shutil.rmtree(profile, ignore_errors=True)

def _calibre_fixture(source: Path, target_extension: str, root: Path) -> Path | None:
    executable = command_path('calibre')
    if not executable:
        return None
    target = root / f'sample{target_extension}'
    completed = subprocess.run([executable, str(source), str(target)], stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=180, check=False)
    return target if completed.returncode == 0 and target.exists() and (target.stat().st_size > 32) else None

def fixtures(root: Path) -> dict[str, Path]:
    files: dict[str, Path] = {}
    pdf = root / 'sample.pdf'
    pdf_canvas = canvas.Canvas(str(pdf))
    pdf_canvas.drawString(72, 760, 'AJN PDF acceptance test 123')
    pdf_canvas.drawString(72, 730, 'English and conversion sample')
    pdf_canvas.save()
    files['.pdf'] = pdf
    image = Image.new('RGB', (1400, 700), 'white')
    draw = ImageDraw.Draw(image)
    font = None
    for candidate in ('C:/Windows/Fonts/arial.ttf', '/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf'):
        try:
            font = ImageFont.truetype(candidate, 72)
            break
        except OSError:
            continue
    draw.text((70, 140), 'AJN PDF TEST 123', fill='black', font=font)
    draw.text((70, 260), 'DOCUMENT CONVERSION', fill='black', font=font)
    for extension, image_format in (('.png', 'PNG'), ('.jpg', 'JPEG'), ('.jpeg', 'JPEG'), ('.webp', 'WEBP'), ('.bmp', 'BMP'), ('.gif', 'GIF'), ('.tif', 'TIFF'), ('.tiff', 'TIFF')):
        destination = root / f'sample{extension}'
        image.save(destination, image_format)
        files[extension] = destination
    image.close()
    svg = root / 'sample.svg'
    svg.write_text('<svg xmlns="http://www.w3.org/2000/svg" width="400" height="200"><rect width="400" height="200" fill="white"/><text x="20" y="80" font-size="28">AJN PDF SVG</text></svg>', encoding='utf-8')
    files['.svg'] = svg
    text_file = root / 'sample.txt'
    text_file.write_text('AJN PDF text conversion sample.\nSecond line.', encoding='utf-8')
    files['.txt'] = text_file
    markdown = root / 'sample.md'
    markdown.write_text('# AJN PDF Markdown conversion sample.', encoding='utf-8')
    files['.md'] = markdown
    files['.markdown'] = markdown
    html = root / 'sample.html'
    html.write_text('<!doctype html><html><body><h1>AJN PDF</h1><p>HTML sample.</p></body></html>', encoding='utf-8')
    files['.html'] = html
    files['.htm'] = html
    xml = root / 'sample.xml'
    xml.write_text('<root><title>AJN PDF</title><value>123</value></root>', encoding='utf-8')
    files['.xml'] = xml
    json_file = root / 'sample.json'
    json_file.write_text(json.dumps({'product': 'AJN PDF', 'value': 123}), encoding='utf-8')
    files['.json'] = json_file
    csv_file = root / 'sample.csv'
    csv_file.write_text('name,value\nAJN PDF,123', encoding='utf-8')
    files['.csv'] = csv_file
    rtf = root / 'sample.rtf'
    rtf.write_text('{\\rtf1\\ansi AJN PDF RTF sample}', encoding='utf-8')
    files['.rtf'] = rtf
    docx = root / 'sample.docx'
    document = Document()
    document.add_heading('AJN PDF', 0)
    document.add_paragraph('DOCX conversion sample.')
    document.save(docx)
    files['.docx'] = docx
    xlsx = root / 'sample.xlsx'
    workbook = Workbook()
    sheet = workbook.active
    sheet.append(['AJN PDF', 'Value'])
    sheet.append(['Sample', 123])
    workbook.save(xlsx)
    files['.xlsx'] = xlsx
    pptx = root / 'sample.pptx'
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = 'AJN PDF'
    slide.placeholders[1].text = 'Presentation sample'
    presentation.save(pptx)
    files['.pptx'] = pptx
    office_variants = [(docx, '.doc', 'doc:MS Word 97'), (docx, '.odt', 'odt'), (xlsx, '.xls', 'xls:MS Excel 97'), (xlsx, '.ods', 'ods'), (pptx, '.ppt', 'ppt:MS PowerPoint 97'), (pptx, '.odp', 'odp')]
    for office_source, extension, filter_name in office_variants:
        generated = _libreoffice_fixture(office_source, extension, filter_name, root)
        if generated:
            files[extension] = generated
    epub_file = root / 'sample.epub'
    book = epub.EpubBook()
    book.set_identifier('ajn-pdf-acceptance')
    book.set_title('AJN PDF Acceptance')
    book.set_language('en')
    chapter = epub.EpubHtml(title='Acceptance', file_name='acceptance.xhtml', lang='en')
    chapter.content = '<h1>AJN PDF</h1><p>eBook conversion acceptance sample.</p>'
    book.add_item(chapter)
    book.toc = (epub.Link('acceptance.xhtml', 'Acceptance', 'acceptance'),)
    book.spine = ['nav', chapter]
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    epub.write_epub(str(epub_file), book)
    files['.epub'] = epub_file
    for extension in ('.mobi', '.azw3'):
        generated = _calibre_fixture(epub_file, extension, root)
        if generated:
            files[extension] = generated
    try:
        import pillow_heif
        pillow_heif.register_heif_opener()
        with Image.open(files['.png']) as heif_source:
            heif_source = heif_source.convert('RGB')
            for extension in ('.heic', '.heif'):
                destination = root / f'sample{extension}'
                heif_source.save(destination, 'HEIF')
                if destination.exists() and destination.stat().st_size > 32:
                    files[extension] = destination
    except Exception:
        pass
    eml = root / 'sample.eml'
    eml.write_text('From: sender@example.com\nTo: receiver@example.com\nSubject: AJN PDF\nContent-Type: text/plain; charset=utf-8 Email conversion sample.', encoding='utf-8')
    files['.eml'] = eml
    return files

def validate_output(target: Path, expected_extension: str) -> None:
    if not target.exists() or target.stat().st_size < 20:
        raise RuntimeError('The generated output is empty.')
    if expected_extension and target.suffix.lower() != expected_extension.lower():
        raise RuntimeError(f'Unexpected output extension:{target.suffix}')
    extension = target.suffix.lower()
    signature = target.read_bytes()[:16]
    if extension == '.pdf':
        if not signature.startswith(b'%PDF-'):
            raise RuntimeError('The generated.pdf does not have a PDF signature.')
        with fitz.open(target) as result:
            if result.page_count < 1:
                raise RuntimeError('The generated PDF contains no pages.')
    elif extension in {'.zip', '.docx', '.xlsx', '.pptx', '.odt', '.ods', '.odp', '.epub'}:
        import zipfile
        if not zipfile.is_zipfile(target):
            raise RuntimeError(f'The generated{extension}is not a valid ZIP-based container.')
        with zipfile.ZipFile(target) as archive:
            if not archive.namelist():
                raise RuntimeError(f'The generated{extension}container is empty.')
    elif extension in {'.png', '.jpg', '.jpeg', '.webp', '.bmp', '.gif', '.tif', '.tiff', '.avif', '.heic', '.heif'}:
        with Image.open(target) as image:
            image.verify()

def validate_workdir_cleanup(workdir: Path) -> None:
    leftovers = []
    for pattern in ('recognized_text-source-*', 'recognized_text-page-*', 'tess-page-*', 'libreoffice-profile-*'):
        leftovers.extend((path for path in workdir.glob(pattern) if path.exists()))
    if leftovers:
        names = ','.join((path.name for path in leftovers[:8]))
        raise RuntimeError(f'Temporary conversion files were not cleaned:{names}')

def main() -> None:
    with tempfile.TemporaryDirectory(prefix='ajn-full-acceptance-') as temporary:
        root = Path(temporary)
        samples = fixtures(root)
        passed: list[str] = []
        skipped: list[tuple[str, str]] = []
        failed: list[tuple[str, str]] = []
        for tool_id, spec in sorted(SPECS.items()):
            available, reason = tool_available(spec)
            if not available:
                skipped.append((tool_id, reason or 'Required dependency is unavailable.'))
                continue
            if spec.processor == 'url_to_pdf':
                skipped.append((tool_id, 'URL security and deployment-network behavior are tested after deployment.'))
                continue
            source: Path | None = None
            for extension in spec.input_extensions:
                source = samples.get(extension)
                if source:
                    break
            if source is None and spec.input_extensions:
                manual_fixture_ids = {'msg-to-pdf', 'xps-to-pdf'}
                reason = f'No generated fixture is available for{spec.input_extensions[0]}.'
                if tool_id in manual_fixture_ids:
                    skipped.append((tool_id, reason + 'Manual real-file QA is required for this binary format.'))
                    continue
                failed.append((tool_id, reason))
                continue
            workdir = root / f'work-{tool_id}'
            workdir.mkdir(exist_ok=True)
            output = workdir / f'result{spec.output_extension}'
            try:
                inputs = [source] if source else []
                if spec.multi_file and source:
                    inputs = [source, source]
                convert(spec, inputs, output, {'language': 'eng', 'dpi': 100, 'quality': 70, 'grayscale': False}, workdir, None)
                validate_output(output, spec.output_extension)
                validate_workdir_cleanup(workdir)
                passed.append(tool_id)
            except Exception as exc:
                failed.append((tool_id, f'{type(exc).__name__}:{exc}'))
        report = {'passed': passed, 'skipped': [{'id': tool_id, 'reason': reason} for tool_id, reason in skipped], 'failed': [{'id': tool_id, 'reason': reason} for tool_id, reason in failed]}
        report_path = Path(__file__).resolve().parent / 'FULL_ACCEPTANCE_RESULTS.json'
        report_path.write_text(json.dumps(report, indent=2), encoding='utf-8')
        if failed:
            first_failures = ','.join((f'{tool_id}:{reason}' for tool_id, reason in failed[:12]))
            raise RuntimeError(f'Acceptance failures:{first_failures}')
        print(f'PASS: full acceptance suite produced valid outputs for{len(passed)}available tools;{len(skipped)}dependency/fixture-specific tools were documented as skipped.')
if __name__ == '__main__':
    main()
